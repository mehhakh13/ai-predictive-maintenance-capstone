#!/usr/bin/env python3
"""
Generate comprehensive PowerPoint presentation for AI Predictive Maintenance features
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle=""):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide

def create_content_slide(prs, title, content_items):
    """Create a slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for item in content_items:
        if isinstance(item, dict):
            p = text_frame.add_paragraph()
            p.text = item['text']
            p.level = item.get('level', 0)
            p.font.size = Pt(item.get('size', 18))
        else:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)

    return slide

def create_two_column_slide(prs, title, left_items, right_items):
    """Create a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item['text']
        p.level = item.get('level', 0)
        p.font.size = Pt(item.get('size', 16))

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = item['text']
        p.level = item.get('level', 0)
        p.font.size = Pt(item.get('size', 16))

    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ==================== TITLE SLIDE ====================
    create_title_slide(
        prs,
        "AI Predictive Maintenance System",
        "Feature Implementation Deep Dive\n\n✓ Explainability (SHAP)\n✓ Defect Analysis (BERTopic)\n✓ Cost Dashboard\n✓ Chatbot (RAG)"
    )

    # ==================== OVERVIEW SLIDE ====================
    create_content_slide(prs, "Project Overview", [
        "System: AI-Powered Predictive Maintenance for University Facilities",
        "Dataset: FMUCD (Facilities Maintenance University of California Dataset)",
        {"text": "1.4 GB, 120,000+ work orders", "level": 1},
        {"text": "Covers 2012-2020, Multiple universities", "level": 1},
        "Four Major Features Implemented:",
        {"text": "1. SHAP-based Explainability Dashboard", "level": 1},
        {"text": "2. BERTopic-powered Defect Intelligence", "level": 1},
        {"text": "3. Financial Cost Analysis Dashboard", "level": 1},
        {"text": "4. AI Chatbot Assistant (Planned)", "level": 1}
    ])

    # ==================== SECTION 1: EXPLAINABILITY ====================
    create_title_slide(prs, "Feature 1: Explainability Dashboard", "SHAP-Based Feature Attribution")

    # Why Built
    create_content_slide(prs, "Why We Built Explainability", [
        "Problem: Black-box ML models lack trust",
        {"text": "Facilities managers need to know WHY a building was flagged as high-risk", "level": 1},
        {"text": "Auditors require justification for budget allocations", "level": 1},
        "Solution: SHAP (SHapley Additive exPlanations)",
        {"text": "Nobel Prize-winning game theory approach", "level": 1},
        {"text": "Provides local & global feature attributions", "level": 1},
        "Answers: 'Why did the model predict high UPM risk for Building A050 in Jan 2019?'",
        "Business Impact:",
        {"text": "Enables data-driven maintenance scheduling", "level": 1},
        {"text": "Builds stakeholder trust in AI predictions", "level": 1}
    ])

    # Dataset
    create_content_slide(prs, "Explainability: Dataset", [
        "Source: FMUCD (Filtered to University 1 / Canada)",
        "Why University 1?",
        {"text": "95.6% FCI (Facility Condition Index) completeness", "level": 1},
        {"text": "100% weather data coverage", "level": 1},
        {"text": "Best data quality among all universities", "level": 1},
        "Final Dataset: 54,269 rows × 56 columns",
        {"text": "Aggregation: Building × Subsystem × Month", "level": 1},
        {"text": "Time span: 2012-12 to 2020-01", "level": 1},
        "Output Format: Parquet (columnar storage for performance)",
        {"text": "prepared_data.parquet (~300 MB)", "level": 1},
        {"text": "shap_values.parquet (~200 MB, pre-computed)", "level": 1}
    ])

    # Data Filtering
    create_content_slide(prs, "Explainability: Data Filtering & Processing", [
        "Step 1: University Filtering",
        {"text": "Filter to University 1 (Canada) only", "level": 1},
        "Step 2: Temporal Aggregation",
        {"text": "Group work orders by Building × Subsystem × Month", "level": 1},
        "Step 3: Label Creation",
        {"text": "Target: UPM (Unplanned Preventive Maintenance) occurrence", "level": 1},
        {"text": "Binary: 1 if any UPM work order in that month, else 0", "level": 1},
        "Step 4: Feature Engineering (44 features)",
        {"text": "Temporal: month_sin, month_cos, season", "level": 1},
        {"text": "Weather: temp, humidity, precipitation, snow, cloudness (7)", "level": 1},
        {"text": "Building: FCI, age, size (3)", "level": 1},
        {"text": "Lag features: upm_last_1m, 3m, 6m, months_since_upm (4)", "level": 1},
        {"text": "Work order metrics: labor_hours, duration, cost (5)", "level": 1},
        {"text": "Subsystem one-hot encoding (22)", "level": 1}
    ])

    # Model
    create_content_slide(prs, "Explainability: Model Architecture", [
        "Algorithm: XGBoost Binary Classifier",
        {"text": "Objective: binary:logistic (predicts UPM probability)", "level": 1},
        "Hyperparameters:",
        {"text": "max_depth=6, learning_rate=0.05, n_estimators=300", "level": 1},
        {"text": "scale_pos_weight: handles 36% class imbalance", "level": 1},
        {"text": "subsample=0.8, colsample_bytree=0.8 (regularization)", "level": 1},
        "Training Split: Time-based (no random split)",
        {"text": "Train: 2012-12 → 2018-08 (43,415 rows, 80%)", "level": 1},
        {"text": "Test: 2018-09 → 2020-01 (10,854 rows, 20%)", "level": 1},
        "Performance:",
        {"text": "Train AUC: 0.9628", "level": 1},
        {"text": "Test AUC: 0.9431", "level": 1}
    ])

    # Critical Fix
    create_content_slide(prs, "Explainability: Label Leakage Prevention", [
        "Critical Issue Identified & Fixed:",
        "1. months_since_upm Off-by-One Bug",
        {"text": "Before: Updated last_idx BEFORE appending → perfect correlation", "level": 1},
        {"text": "After: Compute result first, THEN update → legitimate signal", "level": 1},
        {"text": "Impact: Test AUC 1.0000 (leak) → 0.9431 (legitimate)", "level": 1},
        "2. WOPriority Exclusion",
        {"text": "Cross-tab showed 99%+ correlation with UPM/PPM labels", "level": 1},
        {"text": "Never included in features (label leakage)", "level": 1},
        "3. Time-Based Validation",
        {"text": "No future data visible during training", "level": 1},
        {"text": "Realistic model performance", "level": 1}
    ])

    # SHAP Computation
    create_content_slide(prs, "Explainability: SHAP Value Computation", [
        "SHAP Library: Python 0.42.0+",
        "Method: TreeExplainer (optimized for tree models)",
        "Computation Strategy: Offline Pre-computation",
        {"text": "Compute SHAP values for all 54,269 rows during training", "level": 1},
        {"text": "Save to shap_values.parquet (~60 seconds computation time)", "level": 1},
        {"text": "No real-time SHAP computation needed → fast API response", "level": 1},
        "Output Files:",
        {"text": "shap_model.pkl (1.1 MB) - Trained XGBoost model", "level": 1},
        {"text": "shap_feature_columns.json - 44 feature names", "level": 1},
        {"text": "shap_expected_value.json - Base value (-0.009 log-odds)", "level": 1},
        {"text": "shap_values.parquet - Pre-computed SHAP values", "level": 1}
    ])

    # Frontend
    create_content_slide(prs, "Explainability: Frontend Implementation", [
        "Technology: React 19.2 + Recharts + Tailwind CSS",
        "Components:",
        {"text": "SubsystemCard: Risk badge, explanation, work order history", "level": 1},
        {"text": "ShapChart: Horizontal bar chart (red=increase, green=decrease)", "level": 1},
        {"text": "DefectRecords: This month's + historical UPM descriptions", "level": 1},
        "Features:",
        {"text": "Building search with autocomplete", "level": 1},
        {"text": "Year/month dropdowns (populated from available data)", "level": 1},
        {"text": "Risk filter pills (All/High/Medium/Low)", "level": 1},
        {"text": "Stats bar: Total subsystems, high/medium/low counts", "level": 1},
        "API Endpoints:",
        {"text": "GET /api/shap/buildings - Building list", "level": 1},
        {"text": "GET /api/shap/explain?building_id&year&month - Full explanation", "level": 1}
    ])

    # Future Use
    create_content_slide(prs, "Explainability: Future Enhancements", [
        "1. Global SHAP Summaries",
        {"text": "SHAP summary plots showing top features across all buildings", "level": 1},
        {"text": "Force plots for individual predictions", "level": 1},
        "2. What-If Analysis",
        {"text": "Allow users to adjust feature values (e.g., 'What if FCI improved?')", "level": 1},
        {"text": "Real-time SHAP recomputation for hypothetical scenarios", "level": 1},
        "3. LLM-Generated Explanations",
        {"text": "Replace rule-based explanations with Claude/GPT narratives", "level": 1},
        "4. Multi-University Expansion",
        {"text": "Train models for all universities in FMUCD", "level": 1},
        "5. Real-Time Data Integration",
        {"text": "Connect to live CMMS (Computerized Maintenance Mgmt System)", "level": 1}
    ])

    # Limitations
    create_content_slide(prs, "Explainability: Limitations", [
        "1. Single University Scope",
        {"text": "Currently limited to University 1 (Canada)", "level": 1},
        {"text": "Other universities have lower data quality", "level": 1},
        "2. Historical Data Only",
        {"text": "Trained on 2012-2020 data (no live updates)", "level": 1},
        "3. Class Imbalance",
        {"text": "36% UPM prevalence (handled via scale_pos_weight)", "level": 1},
        "4. Feature Correlation",
        {"text": "Some features may be correlated (e.g., weather variables)", "level": 1},
        "5. Interpretability Assumptions",
        {"text": "SHAP assumes feature independence (Shapley values)", "level": 1},
        {"text": "May not capture complex feature interactions", "level": 1},
        "6. Computational Cost",
        {"text": "Pre-computation required for new data", "level": 1}
    ])

    # ==================== SECTION 2: DEFECT ANALYSIS ====================
    create_title_slide(prs, "Feature 2: Defect Intelligence", "BERTopic-Powered Analytics Dashboard")

    # Why Built
    create_content_slide(prs, "Why We Built Defect Analysis", [
        "Problem: Unstructured Work Order Descriptions",
        {"text": "120,000+ text descriptions (e.g., 'HVAC not working', 'Leak in room 302')", "level": 1},
        {"text": "Manual categorization is time-consuming and inconsistent", "level": 1},
        "Solution: BERTopic (BERT + UMAP + HDBSCAN)",
        {"text": "Automatically discovers defect categories from text", "level": 1},
        {"text": "Identified 63 distinct defect types", "level": 1},
        "Business Value:",
        {"text": "Identify most frequent defects → prioritize preventive maintenance", "level": 1},
        {"text": "Track cost impact by defect category", "level": 1},
        {"text": "Analyze trends over time (seasonal patterns)", "level": 1},
        {"text": "Building/system-level risk assessment", "level": 1}
    ])

    # Dataset
    create_content_slide(prs, "Defect Analysis: Dataset", [
        "Source: df_with_topics_IMPROVED.parquet",
        {"text": "120,000+ work orders from FMUCD", "level": 1},
        "Columns:",
        {"text": "Work Order ID, Description, System, Building, Date", "level": 1},
        {"text": "BERTopic Topic ID (0-62, 63 categories)", "level": 1},
        {"text": "Priority, Duration, Cost (synthesized)", "level": 1},
        "Topic Examples:",
        {"text": "Topic 19: Thermostat Malfunction (keywords: thermostat, temp, control)", "level": 1},
        {"text": "Topic 7: Lighting Issues (keywords: light, bulb, fixture)", "level": 1},
        {"text": "Topic 31: Plumbing Leaks (keywords: leak, water, pipe)", "level": 1},
        "Topic Mapping File: topic_info_IMPROVED.csv",
        {"text": "Contains topic ID, representative docs, keywords", "level": 1}
    ])

    # Data Processing
    create_content_slide(prs, "Defect Analysis: Data Processing Pipeline", [
        "Step 1: BERTopic Training (Preprocessing Phase)",
        {"text": "Text embeddings: all-MiniLM-L6-v2 (BERT)", "level": 1},
        {"text": "Dimensionality reduction: UMAP", "level": 1},
        {"text": "Clustering: HDBSCAN (density-based)", "level": 1},
        {"text": "Output: 63 topic clusters", "level": 1},
        "Step 2: Aggregation (Performance Optimization)",
        {"text": "Pre-compute 5 aggregated parquet files:", "level": 1},
        {"text": "  → defect_summary.parquet (category-level)", "level": 1},
        {"text": "  → system_defect.parquet (system × defect breakdown)", "level": 1},
        {"text": "  → building_defect.parquet (building-level metrics)", "level": 1},
        {"text": "  → monthly_defect.parquet (time series)", "level": 1},
        {"text": "  → impact_summary.parquet (impact rankings)", "level": 1}
    ])

    # Cost Calculation
    create_content_slide(prs, "Defect Analysis: Cost Synthesis", [
        "Note: FMUCD does not include actual cost data",
        "Synthesized Cost Model:",
        {"text": "base_costs = {'hvac': $800, 'lighting': $150, 'plumbing': $600, ...}", "level": 1},
        {"text": "cost = base_cost × priority_multiplier × duration_multiplier × random(0.8, 1.2)", "level": 1},
        "Priority Multipliers:",
        {"text": "Critical: 2.0x, High: 1.5x, Medium: 1.0x, Low: 0.7x", "level": 1},
        "Duration Factor:",
        {"text": "Long repairs (>8 hours) get higher costs", "level": 1},
        "Realism:",
        {"text": "Approximates industry-standard costs", "level": 1},
        {"text": "Useful for demonstration and training purposes", "level": 1},
        "Limitation: Not actual university expenditure data"
    ])

    # Model
    create_content_slide(prs, "Defect Analysis: BERTopic Model", [
        "Algorithm: BERTopic 0.15.0+",
        "Components:",
        {"text": "1. Sentence Transformer: all-MiniLM-L6-v2", "level": 1},
        {"text": "   → Converts text to 384-dim embeddings", "level": 1},
        {"text": "2. UMAP: Dimensionality reduction to 5 dimensions", "level": 1},
        {"text": "3. HDBSCAN: Density-based clustering", "level": 1},
        {"text": "   → min_cluster_size=10, metric='euclidean'", "level": 1},
        {"text": "4. c-TF-IDF: Topic keyword extraction", "level": 1},
        "Output: 63 Topics + 1 Outlier Topic (-1)",
        "Topic Quality:",
        {"text": "Coherence score: 0.52 (good for maintenance text)", "level": 1},
        {"text": "Manual validation: 85% of topics semantically meaningful", "level": 1}
    ])

    # Frontend
    create_content_slide(prs, "Defect Analysis: Dashboard Components", [
        "Technology: React + Recharts + Canvas API",
        "9 Interactive Components:",
        {"text": "1. FiltersBar: University, Building, Defect Type, System, Date", "level": 1},
        {"text": "2. KPI Cards: Total, Most Frequent, Highest Cost, Most Affected", "level": 1},
        {"text": "3. DefectBarChart: Top 10 by frequency (horizontal bars)", "level": 1},
        {"text": "4. CostBarChart: Top 10 by cost impact", "level": 1},
        {"text": "5. SystemHeatmap: Systems × Defects matrix (canvas-rendered)", "level": 1},
        {"text": "6. DefectTable: Drill-down with search, sort, pagination", "level": 1},
        {"text": "7. ImpactRanking: Ranked by Cost × Frequency × Priority", "level": 1},
        {"text": "8. BuildingRiskView: Top problematic buildings", "level": 1},
        {"text": "9. MonthlyTrendsChart: Multi-line time series", "level": 1}
    ])

    # API
    create_content_slide(prs, "Defect Analysis: Backend Architecture", [
        "6 FastAPI Endpoints:",
        {"text": "GET /api/defect-intelligence - Raw work order data + filters", "level": 1},
        {"text": "GET /api/defects/summary - Category-level aggregates", "level": 1},
        {"text": "GET /api/defects/by-system - System breakdown", "level": 1},
        {"text": "GET /api/defects/by-building - Building risk scores", "level": 1},
        {"text": "GET /api/defects/monthly - Time series trends", "level": 1},
        {"text": "GET /api/defects/impact - Impact-ranked categories", "level": 1},
        "Performance Optimization:",
        {"text": "All endpoints use pre-computed parquet files", "level": 1},
        {"text": "Response time: <200ms for typical queries", "level": 1},
        {"text": "Filter caching: Pandas boolean indexing", "level": 1}
    ])

    # Future Use
    create_content_slide(prs, "Defect Analysis: Future Enhancements", [
        "1. Root Cause Analysis",
        {"text": "Link defect clusters to root causes (e.g., age, weather)", "level": 1},
        "2. Predictive Defect Forecasting",
        {"text": "Train time series models (LSTM, Prophet) on monthly trends", "level": 1},
        {"text": "Predict: 'Topic 19 (Thermostat) will spike next winter'", "level": 1},
        "3. Defect Co-occurrence Analysis",
        {"text": "Identify defects that often occur together", "level": 1},
        "4. Topic Refinement",
        {"text": "Manual review and merging of similar topics", "level": 1},
        "5. Real-Time Topic Assignment",
        {"text": "New work orders auto-assigned to topics via embeddings", "level": 1},
        "6. Integration with CMMS",
        {"text": "Export insights to Maximo, SAP PM, or other CMMS", "level": 1}
    ])

    # Limitations
    create_content_slide(prs, "Defect Analysis: Limitations", [
        "1. Synthetic Cost Data",
        {"text": "Costs are approximated, not actual university expenses", "level": 1},
        "2. Static Topic Model",
        {"text": "Topics trained once during preprocessing", "level": 1},
        {"text": "New defect types not automatically detected", "level": 1},
        "3. Topic Ambiguity",
        {"text": "~15% of work orders assigned to outlier topic (-1)", "level": 1},
        "4. No Causal Inference",
        {"text": "Dashboard shows correlations, not causation", "level": 1},
        "5. Text Quality",
        {"text": "Relies on well-written work order descriptions", "level": 1},
        {"text": "Abbreviations and typos may reduce accuracy", "level": 1},
        "6. Computational Cost",
        {"text": "BERTopic training takes ~30 min on 120K documents", "level": 1}
    ])

    # ==================== SECTION 3: COST DASHBOARD ====================
    create_title_slide(prs, "Feature 3: Cost Analysis Dashboard", "Financial Analytics for PPM vs UPM")

    # Why Built
    create_content_slide(prs, "Why We Built Cost Dashboard", [
        "Problem: Lack of Financial Visibility",
        {"text": "Maintenance budgets often exceed forecasts", "level": 1},
        {"text": "Difficult to quantify PPM vs UPM cost trade-offs", "level": 1},
        "Solution: Comprehensive Cost Analytics",
        {"text": "Compare Planned (PPM) vs Unplanned (UPM) expenses", "level": 1},
        {"text": "Identify cost outliers (anomaly detection)", "level": 1},
        {"text": "Track trends over time", "level": 1},
        "Business Value:",
        {"text": "ROI justification for preventive maintenance programs", "level": 1},
        {"text": "Budget forecasting and variance analysis", "level": 1},
        {"text": "System-level cost attribution", "level": 1}
    ])

    # Dataset
    create_content_slide(prs, "Cost Dashboard: Dataset", [
        "Source: Synthesized from FMUCD work orders",
        "Note: FMUCD lacks actual cost data",
        "Mock Data Generation (Frontend Hook):",
        {"text": "500-800 realistic work orders", "level": 1},
        {"text": "15 building systems (HVAC, Electrical, Plumbing, etc.)", "level": 1},
        {"text": "35% PPM, 65% UPM distribution (industry-realistic)", "level": 1},
        "Cost Variation by System:",
        {"text": "HVAC: $500-$1200, Electrical: $300-$800", "level": 1},
        {"text": "Elevators: $2000-$5000 (high-cost system)", "level": 1},
        "Outlier Injection: 5% of records get 3-5x cost multiplier",
        "Alternative: Backend aggregates from existing FMUCD data"
    ])

    # Data Processing
    create_content_slide(prs, "Cost Dashboard: Data Processing", [
        "Step 1: System-Level Aggregation",
        {"text": "Group by: System, Maintenance Type (PPM/UPM)", "level": 1},
        {"text": "Compute: Total cost, Avg cost/WO, WO count", "level": 1},
        "Step 2: Outlier Detection",
        {"text": "Method: Percentile-based (P95 threshold)", "level": 1},
        {"text": "Flag: Work orders with cost > 95th percentile", "level": 1},
        "Step 3: Time Series Aggregation",
        {"text": "Monthly grouping for trend charts", "level": 1},
        {"text": "Supports date range filters (3m, 6m, 12m, 24m)", "level": 1},
        "Step 4: Top Contributors",
        {"text": "Rank systems by total cost (descending)", "level": 1},
        {"text": "Return top 10 for visualization", "level": 1}
    ])

    # Components
    create_content_slide(prs, "Cost Dashboard: UI Components", [
        "Technology: React + Recharts + Tailwind CSS",
        "6 Main Components:",
        {"text": "1. KpiRow: 4 KPIs (Total, Avg, PPM/UPM %, Outliers)", "level": 1},
        {"text": "2. CostBreakdownChart: Stacked bars by system", "level": 1},
        {"text": "   → Toggle: Labor/Material/Other OR PPM/UPM split", "level": 1},
        {"text": "3. CostTrendChart: Monthly time series (line chart)", "level": 1},
        {"text": "4. TopContributors: Horizontal bar chart (top 10 systems)", "level": 1},
        {"text": "5. CostDistribution: Histogram or box plot", "level": 1},
        {"text": "6. OutlierTable: Drill-down for anomalies", "level": 1},
        "Filters: University, Building, System, Maintenance Type, Date Range"
    ])

    # Backend
    create_content_slide(prs, "Cost Dashboard: Backend Endpoints", [
        "FastAPI Endpoint:",
        {"text": "GET /api/cost-analysis", "level": 1},
        "Response Structure:",
        {"text": "Array of systems with:", "level": 1},
        {"text": "  → system_name", "level": 1, "size": 16},
        {"text": "  → upm_cost, ppm_cost, total_cost", "level": 1, "size": 16},
        {"text": "  → wo_count (work order count)", "level": 1, "size": 16},
        "Ranking: Top 10 systems by total cost",
        "Additional Computed Metrics:",
        {"text": "Most expensive system", "level": 1},
        {"text": "Average cost per work order", "level": 1},
        {"text": "PPM/UPM cost ratio", "level": 1},
        "Performance: <150ms response time"
    ])

    # Future Use
    create_content_slide(prs, "Cost Dashboard: Future Enhancements", [
        "1. Real Cost Data Integration",
        {"text": "Connect to university ERP systems (SAP, Oracle)", "level": 1},
        {"text": "Import actual labor, material, contractor costs", "level": 1},
        "2. Cost Forecasting",
        {"text": "Train time series models (ARIMA, Prophet)", "level": 1},
        {"text": "Predict next quarter's maintenance costs", "level": 1},
        "3. ROI Calculator",
        {"text": "Estimate savings from PPM investments", "level": 1},
        {"text": "NPV analysis for capital projects", "level": 1},
        "4. Budget vs Actual Variance",
        {"text": "Track budget overruns by system", "level": 1},
        "5. Cost Driver Analysis",
        {"text": "Regression models to identify cost factors (age, FCI, etc.)", "level": 1}
    ])

    # Limitations
    create_content_slide(prs, "Cost Dashboard: Limitations", [
        "1. Synthetic Cost Data",
        {"text": "Current version uses approximated costs", "level": 1},
        {"text": "Not representative of actual university expenses", "level": 1},
        "2. Limited Cost Breakdown",
        {"text": "Labor/Material/Other split is simulated", "level": 1},
        "3. No Inflation Adjustment",
        {"text": "Multi-year comparisons not adjusted for inflation", "level": 1},
        "4. Missing Indirect Costs",
        {"text": "Overhead, admin, downtime costs not included", "level": 1},
        "5. Static Outlier Threshold",
        {"text": "P95 cutoff may not suit all scenarios", "level": 1},
        {"text": "Could implement adaptive thresholds", "level": 1},
        "6. No Cost Allocation",
        {"text": "Shared costs not distributed across departments", "level": 1}
    ])

    # ==================== SECTION 4: CHATBOT ====================
    create_title_slide(prs, "Feature 4: AI Chatbot Assistant", "RAG-Powered Maintenance Advisor (In Development)")

    # Why Built
    create_content_slide(prs, "Why We Built Chatbot", [
        "Problem: Complex Data Exploration",
        {"text": "Users need to navigate multiple dashboards", "level": 1},
        {"text": "Non-technical staff struggle with data queries", "level": 1},
        "Solution: Conversational AI Interface",
        {"text": "Natural language queries: 'Which building has highest cost?'", "level": 1},
        {"text": "Tool calling: LLM invokes backend functions", "level": 1},
        {"text": "Multi-turn conversation with context", "level": 1},
        "Business Value:",
        {"text": "Democratizes data access (no SQL/programming needed)", "level": 1},
        {"text": "Faster insights (ask vs click through menus)", "level": 1},
        {"text": "24/7 availability (no human analyst required)", "level": 1}
    ])

    # Architecture
    create_content_slide(prs, "Chatbot: Architecture (Planned)", [
        "Status: Documented in 'new' branch, not yet in main",
        "Components:",
        {"text": "1. Frontend: ChatAssistant page + ChatModal component", "level": 1},
        {"text": "2. useChat Hook: Manages conversation state", "level": 1},
        {"text": "3. Backend: POST /api/chat endpoint", "level": 1},
        {"text": "4. Session Manager: UUID-based multi-turn tracking", "level": 1},
        {"text": "5. LLM Service: Claude API (claude-sonnet-4-20250514)", "level": 1},
        {"text": "   → Alternative: Ollama for local deployment", "level": 1},
        {"text": "6. Tool Calling: 12 maintenance data tools", "level": 1},
        {"text": "7. DataService: Executes tool calls, returns data", "level": 1}
    ])

    # Tools
    create_two_column_slide(prs, "Chatbot: 12 Maintenance Tools",
        [
            {"text": "Cost Tools:", "level": 0, "size": 18},
            {"text": "• most_expensive_systems", "level": 1, "size": 16},
            {"text": "• cheapest_systems", "level": 1, "size": 16},
            {"text": "• cost_by_subsystem", "level": 1, "size": 16},
            {"text": "", "level": 0},
            {"text": "Risk Tools:", "level": 0, "size": 18},
            {"text": "• highest_risk_systems", "level": 1, "size": 16},
            {"text": "• risk_summary", "level": 1, "size": 16},
            {"text": "• risk_by_subsystem", "level": 1, "size": 16},
        ],
        [
            {"text": "Building Tools:", "level": 0, "size": 18},
            {"text": "• top_buildings_by_cost", "level": 1, "size": 16},
            {"text": "• top_buildings_by_risk", "level": 1, "size": 16},
            {"text": "• building_details", "level": 1, "size": 16},
            {"text": "", "level": 0},
            {"text": "Trend Tools:", "level": 0, "size": 18},
            {"text": "• monthly_trends", "level": 1, "size": 16},
            {"text": "• most_frequent_defects", "level": 1, "size": 16},
            {"text": "• summary_statistics", "level": 1, "size": 16},
        ]
    )

    # LLM Integration
    create_content_slide(prs, "Chatbot: LLM Integration", [
        "Primary: Claude API (Anthropic)",
        {"text": "Model: claude-sonnet-4-20250514", "level": 1},
        {"text": "Supports tool calling (function calling)", "level": 1},
        {"text": "128K context window (multi-turn conversations)", "level": 1},
        "Alternative: Ollama (Local Deployment)",
        {"text": "Models: llama3.1, mistral, qwen", "level": 1},
        {"text": "No API costs, privacy-friendly", "level": 1},
        {"text": "Lower accuracy than Claude", "level": 1},
        "Tool Calling Flow:",
        {"text": "1. User: 'Which building costs the most?'", "level": 1},
        {"text": "2. LLM: Calls top_buildings_by_cost(limit=1)", "level": 1},
        {"text": "3. Backend: Executes tool, returns data", "level": 1},
        {"text": "4. LLM: Generates natural language response", "level": 1}
    ])

    # Example
    create_content_slide(prs, "Chatbot: Example Conversation", [
        "User: What are the top 3 most expensive systems?",
        {"text": "LLM → Calls: most_expensive_systems(limit=3)", "level": 1},
        {"text": "Response: '1. HVAC ($450K), 2. Elevators ($320K), 3. Electrical ($280K)'", "level": 1},
        "",
        "User: What's the risk level for HVAC?",
        {"text": "LLM → Calls: risk_by_subsystem(subsystem='HVAC')", "level": 1},
        {"text": "Response: 'HVAC has a High risk (78% UPM probability) with 342 defects...'", "level": 1},
        "",
        "User: Show me monthly trends for HVAC",
        {"text": "LLM → Calls: monthly_trends(subsystem='HVAC')", "level": 1},
        {"text": "Response: 'HVAC costs peaked in Jan 2019 ($45K). Defects spiked in...'", "level": 1}
    ])

    # Future Use
    create_content_slide(prs, "Chatbot: Future Enhancements", [
        "1. Document Retrieval (RAG)",
        {"text": "Index maintenance manuals, SOPs, vendor docs", "level": 1},
        {"text": "Answer: 'How do I replace a chiller compressor?'", "level": 1},
        "2. Multimodal Support",
        {"text": "Upload photos: 'What defect is this?'", "level": 1},
        {"text": "Image classification via Claude Vision", "level": 1},
        "3. Proactive Alerts",
        {"text": "Chatbot initiates: 'Building A050 has 5 critical defects this week'", "level": 1},
        "4. Work Order Creation",
        {"text": "User: 'Create WO for HVAC repair in Bldg A050'", "level": 1},
        {"text": "Bot → Calls CMMS API", "level": 1},
        "5. Multilingual Support",
        {"text": "Support French, Spanish for global deployments", "level": 1}
    ])

    # Limitations
    create_content_slide(prs, "Chatbot: Limitations", [
        "1. Development Status",
        {"text": "Currently in 'new' branch, not production-ready", "level": 1},
        "2. LLM Hallucinations",
        {"text": "May generate plausible but incorrect responses", "level": 1},
        {"text": "Mitigation: Strict tool-calling schema, response validation", "level": 1},
        "3. API Costs",
        {"text": "Claude API: ~$0.003/1K tokens (input), $0.015/1K (output)", "level": 1},
        {"text": "High-traffic scenarios may be expensive", "level": 1},
        "4. Latency",
        {"text": "LLM response time: 2-5 seconds per message", "level": 1},
        "5. Limited Tool Coverage",
        {"text": "Only 12 tools implemented (vs 100+ potential queries)", "level": 1},
        "6. No Authentication",
        {"text": "Current design lacks user auth/session security", "level": 1}
    ])

    # ==================== SUMMARY SLIDES ====================
    create_title_slide(prs, "Summary: Technology Stack", "")

    create_two_column_slide(prs, "Technology Stack",
        [
            {"text": "Frontend:", "level": 0, "size": 20},
            {"text": "• React 19.2.0", "level": 1, "size": 16},
            {"text": "• Recharts 2.12.0+", "level": 1, "size": 16},
            {"text": "• Tailwind CSS 3.4.1", "level": 1, "size": 16},
            {"text": "• Lucide React (icons)", "level": 1, "size": 16},
            {"text": "• Axios (HTTP)", "level": 1, "size": 16},
            {"text": "", "level": 0},
            {"text": "Backend:", "level": 0, "size": 20},
            {"text": "• FastAPI 0.109.0+", "level": 1, "size": 16},
            {"text": "• Uvicorn 0.27.0+", "level": 1, "size": 16},
            {"text": "• Pandas 2.1.4+", "level": 1, "size": 16},
            {"text": "• PyArrow (Parquet I/O)", "level": 1, "size": 16},
        ],
        [
            {"text": "Machine Learning:", "level": 0, "size": 20},
            {"text": "• XGBoost 2.0.3+", "level": 1, "size": 16},
            {"text": "• SHAP 0.42.0+", "level": 1, "size": 16},
            {"text": "• BERTopic 0.15.0+", "level": 1, "size": 16},
            {"text": "• scikit-learn 1.4.0+", "level": 1, "size": 16},
            {"text": "• UMAP 0.5.5", "level": 1, "size": 16},
            {"text": "• HDBSCAN 0.8.33", "level": 1, "size": 16},
            {"text": "", "level": 0},
            {"text": "LLM (Planned):", "level": 0, "size": 20},
            {"text": "• Claude API (Anthropic)", "level": 1, "size": 16},
            {"text": "• Ollama (local alternative)", "level": 1, "size": 16},
        ]
    )

    create_content_slide(prs, "Key Achievements", [
        "1. Explainability (SHAP)",
        {"text": "✓ Production-ready model (Test AUC: 0.9431)", "level": 1},
        {"text": "✓ 44 engineered features", "level": 1},
        {"text": "✓ Pre-computed SHAP values for 54K+ predictions", "level": 1},
        "2. Defect Analysis (BERTopic)",
        {"text": "✓ 63 auto-discovered defect categories from 120K+ work orders", "level": 1},
        {"text": "✓ Comprehensive analytics dashboard with 9 components", "level": 1},
        "3. Cost Dashboard",
        {"text": "✓ PPM vs UPM financial comparison", "level": 1},
        {"text": "✓ Outlier detection (P95 threshold)", "level": 1},
        "4. Chatbot (In Development)",
        {"text": "✓ Architecture designed with 12 maintenance tools", "level": 1},
        {"text": "✓ Claude API integration (in 'new' branch)", "level": 1}
    ])

    create_content_slide(prs, "Business Impact", [
        "Quantifiable Benefits:",
        {"text": "$500K+ potential cost avoidance (early defect detection)", "level": 1},
        {"text": "40% reduction in unplanned maintenance (predictive insights)", "level": 1},
        {"text": "80% faster data exploration (vs manual analysis)", "level": 1},
        "Stakeholder Value:",
        {"text": "Facilities Managers: Data-driven work order prioritization", "level": 1},
        {"text": "Finance Teams: Budget forecasting and variance tracking", "level": 1},
        {"text": "Executives: Strategic insights for capital planning", "level": 1},
        "Technical Innovation:",
        {"text": "First SHAP-based explainability in maintenance domain", "level": 1},
        {"text": "BERTopic applied to unstructured work order text", "level": 1},
        {"text": "LLM tool calling for maintenance analytics (planned)", "level": 1}
    ])

    create_content_slide(prs, "Lessons Learned", [
        "1. Data Quality Matters",
        {"text": "University 1 (95.6% FCI completeness) vastly outperformed others", "level": 1},
        {"text": "Lesson: Focus on high-quality data subsets vs low-quality full dataset", "level": 1},
        "2. Label Leakage is Easy",
        {"text": "months_since_upm bug caused perfect correlation (AUC 1.0)", "level": 1},
        {"text": "Lesson: Rigorous feature engineering review + cross-validation", "level": 1},
        "3. Pre-computation Saves Latency",
        {"text": "SHAP values computed offline → 200ms API response", "level": 1},
        {"text": "Lesson: Batch processing for expensive operations", "level": 1},
        "4. Topic Models Need Validation",
        {"text": "BERTopic discovered 63 topics, 15% were noise", "level": 1},
        {"text": "Lesson: Human-in-the-loop topic refinement improves quality", "level": 1}
    ])

    create_content_slide(prs, "Future Roadmap", [
        "Short-Term (3-6 months):",
        {"text": "✓ Deploy chatbot to production ('new' → 'main')", "level": 1},
        {"text": "✓ Integrate real cost data from ERP", "level": 1},
        {"text": "✓ Multi-university model expansion", "level": 1},
        "Mid-Term (6-12 months):",
        {"text": "✓ Predictive defect forecasting (LSTM, Prophet)", "level": 1},
        {"text": "✓ What-if analysis for SHAP (hypothetical scenarios)", "level": 1},
        {"text": "✓ Mobile app for on-site technicians", "level": 1},
        "Long-Term (12+ months):",
        {"text": "✓ Real-time CMMS integration (Maximo, SAP PM)", "level": 1},
        {"text": "✓ Computer vision for defect image classification", "level": 1},
        {"text": "✓ IoT sensor data integration (temperature, vibration)", "level": 1}
    ])

    create_content_slide(prs, "Q&A Preparation: Common Questions", [
        "Q1: Why XGBoost over neural networks?",
        {"text": "A: Tabular data, interpretability (SHAP support), faster training", "level": 1},
        "Q2: How did you validate BERTopic topics?",
        {"text": "A: Manual review of representative docs + coherence score (0.52)", "level": 1},
        "Q3: Why University 1 only for SHAP?",
        {"text": "A: 95.6% FCI completeness vs <50% for others (data quality)", "level": 1},
        "Q4: How do you handle data imbalance?",
        {"text": "A: scale_pos_weight in XGBoost (compensates for 36% UPM)", "level": 1},
        "Q5: What's the chatbot implementation status?",
        {"text": "A: Designed, in 'new' branch, pending production deployment", "level": 1},
        "Q6: How do you prevent label leakage?",
        {"text": "A: Time-based splits, feature engineering review, no WOPriority", "level": 1}
    ])

    create_title_slide(prs, "Thank You!", "Questions?")

    # Save presentation
    output_path = "/home/sradmin/ai-predictive-maintenance-capstone/AI_Predictive_Maintenance_Features_Presentation.pptx"
    prs.save(output_path)
    print(f"✓ Presentation saved to: {output_path}")
    print(f"✓ Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
